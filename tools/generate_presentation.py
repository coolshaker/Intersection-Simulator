from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "presentation_results.json"
OUTPUT_PATH = ROOT / "Intersection_Simulator_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(17, 47, 78)
TEAL = RGBColor(48, 116, 122)
SKY = RGBColor(93, 162, 190)
ORANGE = RGBColor(210, 135, 55)
INK = RGBColor(36, 46, 56)
MUTED = RGBColor(96, 109, 121)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
ROAD = RGBColor(90, 96, 104)
RED = RGBColor(171, 59, 59)
GREEN = RGBColor(55, 149, 88)


def run_analysis() -> dict:
    subprocess.run(
        ["node", str(ROOT / "tools" / "sim_analysis.js"), str(DATA_PATH)],
        check=True,
        cwd=ROOT,
    )
    return json.loads(DATA_PATH.read_text(encoding="utf-8"))


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(8.8), Inches(0.34))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.47), Inches(0.5), Inches(8.8), Inches(0.22))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(10.5)
        p.font.color.rgb = RGBColor(220, 231, 240)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, left, top, width, height, items, font_size=20, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        first = False
    return box


def add_section_label(slide, text, left, top, width=Inches(2.5)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = SKY
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY


def add_metric_card(slide, left, top, width, title, value, accent):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)

    label = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.16), width - Inches(0.25), Inches(0.25))
    p = label.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    number = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.43), width - Inches(0.25), Inches(0.45))
    p = number.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = accent


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def style_chart(chart, fill_color, value_format="0.0"):
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.maximum_scale = None
    chart.value_axis.format.line.color.rgb = RGBColor(210, 218, 227)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = INK
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = INK
    for series in chart.series:
        series.has_data_labels = True
        series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        series.data_labels.number_format = value_format
        series.data_labels.font.size = Pt(9)
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = fill_color


def add_title_slide(prs: Presentation, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(3.1), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY
    accent.line.fill.background()

    add_text(slide, Inches(0.6), Inches(0.85), Inches(2.0), Inches(0.7), "CEE 370", size=16, bold=True, color=WHITE)
    add_text(slide, Inches(3.55), Inches(1.1), Inches(8.7), Inches(0.9), "Intersection Simulator", size=26, bold=True, color=NAVY)
    add_text(
        slide,
        Inches(3.58),
        Inches(1.95),
        Inches(8.7),
        Inches(0.85),
        "Academic presentation on the problem, methodology, simulation design, results, and future work.",
        size=17,
        color=INK,
    )

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(3.0), Inches(8.95), Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = WHITE
    band.line.color.rgb = SKY
    band.line.width = Pt(1.5)

    add_text(slide, Inches(3.85), Inches(3.28), Inches(8.1), Inches(0.28), "Project context", size=12, bold=True, color=TEAL)
    add_text(
        slide,
        Inches(3.85),
        Inches(3.58),
        Inches(7.9),
        Inches(0.35),
        "Browser-based signalized intersection simulator for teaching traffic operations and control concepts.",
        size=16,
        color=INK,
    )

    add_text(slide, Inches(3.58), Inches(5.3), Inches(4.2), Inches(0.28), "Prepared for", size=11, bold=True, color=MUTED)
    add_text(slide, Inches(3.58), Inches(5.55), Inches(4.2), Inches(0.28), "Transportation Fundamentals", size=15, color=INK)
    add_text(slide, Inches(8.2), Inches(5.3), Inches(4.0), Inches(0.28), "Author", size=11, bold=True, color=MUTED)
    add_text(slide, Inches(8.2), Inches(5.55), Inches(4.0), Inches(0.45), "Dr. Kun Xie, Old Dominion University", size=15, color=INK)
    add_footer(slide, f"{data['project']} | Generated from local project files and reproducible scenario analysis")


def add_problem_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Problem And Motivation", "Why this simulator matters in a transportation fundamentals setting")

    add_section_label(slide, "Core Problem", Inches(0.55), Inches(1.12))
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.5),
        Inches(5.8),
        Inches(4.3),
        [
            "Students often learn signal timing, queueing, and delay from equations alone, without seeing the operational consequences in real time.",
            "A compact classroom simulator can connect traffic engineering theory to observable system behavior at a four-leg signalized intersection.",
            "The project is intended for concept exploration and instruction rather than high-fidelity field calibration.",
        ],
        font_size=18,
    )

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.55), Inches(5.85), Inches(1.5))
    callout.fill.solid()
    callout.fill.fore_color.rgb = LIGHT
    callout.line.color.rgb = SKY
    callout.line.width = Pt(1.5)
    add_text(slide, Inches(7.05), Inches(1.82), Inches(5.3), Inches(0.22), "Guiding Question", size=13, bold=True, color=TEAL)
    add_text(
        slide,
        Inches(7.05),
        Inches(2.1),
        Inches(5.2),
        Inches(0.6),
        "How do signal timing, arrival randomness, turning behavior, and vehicle mix influence queue formation and delay?",
        size=19,
        bold=True,
        color=NAVY,
    )

    add_section_label(slide, "Project Scope", Inches(6.8), Inches(3.35))
    add_bullets(
        slide,
        Inches(6.95),
        Inches(3.72),
        Inches(5.3),
        Inches(2.5),
        [
            "One inbound lane per approach",
            "Two-phase signal control",
            "User-adjustable demand and traffic composition",
            "Live metrics for arrivals, departures, queue, and delay",
        ],
        font_size=17,
    )
    add_footer(slide, "Teaching focus: visualization, experimentation, and interpretation of traffic operations")


def add_methodology_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Methodology", "How the simulator translates engineering assumptions into dynamic system behavior")

    steps = [
        ("1. Define Inputs", "Signal timings, arrival rates by approach, free-flow speed, left-turn ratio, large-vehicle ratio, and simulation speed."),
        ("2. Generate Demand", "Vehicle arrivals are sampled stochastically for each approach, creating natural variability in queue growth."),
        ("3. Update Signals", "The controller alternates east-west and north-south phases with adjustable yellow clearance intervals."),
        ("4. Move Vehicles", "Vehicles follow lane-based trajectories, obey car-following logic, and check permissive left-turn conflicts."),
        ("5. Compute Metrics", "The simulator continuously reports queue counts, average queue, maximum queue, and average delay."),
    ]

    y = 1.4
    for idx, (title, body) in enumerate(steps):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y), Inches(12.0), Inches(0.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE if idx % 2 == 0 else LIGHT
        card.line.color.rgb = SKY
        card.line.width = Pt(1)
        add_text(slide, Inches(1.0), Inches(y + 0.16), Inches(2.25), Inches(0.22), title, size=14, bold=True, color=TEAL)
        add_text(slide, Inches(3.0), Inches(y + 0.13), Inches(9.2), Inches(0.38), body, size=15, color=INK)
        y += 1.0

    add_text(
        slide,
        Inches(0.95),
        Inches(6.65),
        Inches(11.6),
        Inches(0.32),
        "Methodological emphasis: simple enough for classroom use, but rich enough to show the interaction between control, demand, and vehicle behavior.",
        size=13,
        color=MUTED,
    )
    add_footer(slide, "Model class: interactive teaching simulator with stochastic arrivals and simplified operational logic")


def add_design_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Simulation Design", "Geometry, control structure, and performance outputs represented in the application")

    schematic = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.28), Inches(5.1), Inches(4.7))
    schematic.fill.solid()
    schematic.fill.fore_color.rgb = WHITE
    schematic.line.color.rgb = RGBColor(208, 216, 225)
    add_text(slide, Inches(0.88), Inches(1.45), Inches(2.8), Inches(0.25), "Intersection Layout", size=14, bold=True, color=TEAL)

    road_h = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.15), Inches(3.1), Inches(4.0), Inches(0.95))
    road_h.fill.solid()
    road_h.fill.fore_color.rgb = ROAD
    road_h.line.fill.background()
    road_v = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.67), Inches(1.95), Inches(0.95), Inches(3.2))
    road_v.fill.solid()
    road_v.fill.fore_color.rgb = ROAD
    road_v.line.fill.background()

    for x, y, color in [
        (2.4, 2.72, GREEN),
        (3.8, 2.72, RED),
        (2.4, 4.3, RED),
        (3.8, 4.3, GREEN),
    ]:
        sig = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
        sig.fill.solid()
        sig.fill.fore_color.rgb = color
        sig.line.fill.background()

    add_text(slide, Inches(1.12), Inches(2.7), Inches(1.0), Inches(0.2), "EB", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.12), Inches(3.45), Inches(1.0), Inches(0.2), "WB", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.76), Inches(2.05), Inches(0.8), Inches(0.2), "SB", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.74), Inches(4.75), Inches(0.8), Inches(0.2), "NB", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(6.15), Inches(1.46), Inches(2.8), Inches(0.25), "Component Architecture", size=14, bold=True, color=TEAL)
    boxes = [
        (6.15, 1.95, 2.35, 0.78, "User Inputs"),
        (9.05, 1.95, 2.35, 0.78, "Signal Logic"),
        (6.15, 3.15, 2.35, 0.78, "Vehicle Generation"),
        (9.05, 3.15, 2.35, 0.78, "Movement + Yield"),
        (7.6, 4.4, 2.35, 0.78, "Performance Metrics"),
    ]
    for left, top, width, height, label in boxes:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = SKY
        add_text(slide, Inches(left + 0.18), Inches(top + 0.22), Inches(width - 0.3), Inches(0.2), label, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    connectors = [
        ((8.5, 2.34), (9.05, 2.34)),
        ((8.5, 3.54), (9.05, 3.54)),
        ((7.3, 2.73), (7.3, 3.15)),
        ((10.2, 2.73), (10.2, 3.15)),
        ((8.25, 3.93), (8.78, 4.4)),
        ((10.2, 3.93), (9.75, 4.4)),
    ]
    for (x1, y1), (x2, y2) in connectors:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = TEAL
        line.line.width = Pt(1.5)

    add_text(
        slide,
        Inches(6.15),
        Inches(5.55),
        Inches(5.85),
        Inches(0.5),
        "Outputs shown in the interface include arrivals, departures, current queue, average queue, maximum queue, and average delay.",
        size=14,
        color=INK,
    )
    add_footer(slide, "Design intent: expose the major mechanisms of signalized intersection operation without unnecessary software complexity")


def add_results_slide(prs: Presentation, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Results", "Illustrative scenario analysis generated from the simulator logic")

    scenarios = data["scenarios"]
    baseline = scenarios[0]

    add_metric_card(slide, Inches(0.65), Inches(1.22), Inches(2.2), "Baseline Avg Delay", f"{baseline['meanAvgDelaySeconds']} s/veh", ORANGE)
    add_metric_card(slide, Inches(2.95), Inches(1.22), Inches(2.2), "Baseline Avg Queue", f"{baseline['meanAvgQueueVehicles']} veh", TEAL)
    add_metric_card(slide, Inches(5.25), Inches(1.22), Inches(2.2), "Baseline Max Queue", f"{baseline['meanMaxQueueVehicles']} veh", SKY)
    add_metric_card(slide, Inches(7.55), Inches(1.22), Inches(2.2), "Baseline Throughput", f"{baseline['meanThroughputVehPerHour']} veh/h", NAVY)

    delay_data = CategoryChartData()
    delay_data.categories = [item["title"].replace("Scenario ", "") for item in scenarios]
    delay_data.add_series("Average Delay (s/veh)", [item["meanAvgDelaySeconds"] for item in scenarios])
    delay_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.75),
        Inches(2.8),
        Inches(5.7),
        Inches(2.45),
        delay_data,
    ).chart
    delay_chart.value_axis.axis_title.text_frame.text = "s/veh"
    style_chart(delay_chart, TEAL, "0.0")

    queue_data = CategoryChartData()
    queue_data.categories = [item["title"].replace("Scenario ", "") for item in scenarios]
    queue_data.add_series("Maximum Queue (veh)", [item["meanMaxQueueVehicles"] for item in scenarios])
    queue_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(6.85),
        Inches(2.8),
        Inches(5.65),
        Inches(2.45),
        queue_data,
    ).chart
    queue_chart.value_axis.axis_title.text_frame.text = "vehicles"
    style_chart(queue_chart, ORANGE, "0.0")

    rows = len(scenarios) + 1
    table = slide.shapes.add_table(rows, 5, Inches(0.75), Inches(5.55), Inches(11.75), Inches(1.18)).table
    headers = ["Scenario", "Timing (EW-Y-NS)", "Demand (veh/h)", "Avg Delay", "Throughput"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for row_idx, item in enumerate(scenarios, start=1):
        values = [
            item["title"].replace("Scenario ", ""),
            f"{item['timing']['eastWestGreen']}-{item['timing']['yellow']}-{item['timing']['northSouthGreen']}",
            f"EW {item['demandVehPerHour']['west']}/{item['demandVehPerHour']['east']}, NS {item['demandVehPerHour']['north']}/{item['demandVehPerHour']['south']}",
            f"{item['meanAvgDelaySeconds']} s",
            f"{item['meanThroughputVehPerHour']} veh/h",
        ]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)
            p.font.color.rgb = INK
            p.alignment = PP_ALIGN.CENTER

    add_footer(
        slide,
        "Ten replications per scenario over 1,200 simulated seconds. Results are illustrative and intended for instructional interpretation.",
    )


def add_interpretation_slide(prs: Presentation, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Interpretation", "What the scenario results suggest about system behavior")

    scenarios = data["scenarios"]
    scenario_b = scenarios[1]
    scenario_c = scenarios[2]

    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(6.1),
        Inches(4.6),
        [
            "Under balanced demand, the default 30-5-30 timing produces moderate queues and delay, which is appropriate for classroom demonstration of stable operations.",
            "When east-west demand is increased without retiming, delay and queue growth rise noticeably, showing the sensitivity of performance to volume imbalance.",
            f"Reallocating green time to a 40-5-20 split reduces average delay from {scenario_b['meanAvgDelaySeconds']} to {scenario_c['meanAvgDelaySeconds']} s/veh under the heavy east-west case.",
            "The simulator therefore supports comparative reasoning about timing plans, not just passive visualization.",
        ],
        font_size=18,
    )

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.25), Inches(3.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = SKY

    add_text(slide, Inches(7.45), Inches(1.9), Inches(4.8), Inches(0.25), "Key Academic Takeaways", size=14, bold=True, color=TEAL)
    add_bullets(
        slide,
        Inches(7.45),
        Inches(2.25),
        Inches(4.65),
        Inches(2.8),
        [
            "Signal timing affects both queue accumulation and experienced delay.",
            "Random arrivals create run-to-run variation, so replication matters.",
            "Permissive left turns introduce interaction effects beyond simple green-red control.",
            "The tool is best used to teach mechanisms and comparative trends.",
        ],
        font_size=16,
    )
    add_footer(slide, "Interpretation focus: relative differences and operational insight rather than field calibration")


def add_future_work_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide, "Future Work", "Logical extensions for stronger analytical realism and classroom value")

    columns = [
        (
            "Model Enhancements",
            [
                "Add protected left-turn phasing and additional lane groups.",
                "Represent pedestrian calls and more detailed clearance logic.",
                "Introduce approach-specific turning proportions and heavy-vehicle effects.",
            ],
        ),
        (
            "Validation And Analysis",
            [
                "Calibrate selected parameters against observed field data.",
                "Add automated batch experiments, confidence intervals, and exportable reports.",
                "Compare simulated measures with deterministic HCM-style calculations.",
            ],
        ),
        (
            "Instructional Features",
            [
                "Save and reload classroom scenarios for lab exercises.",
                "Add guided experiments that connect inputs to expected hypotheses.",
                "Provide richer visual diagnostics such as phase timelines and queue heatmaps.",
            ],
        ),
    ]

    x_positions = [0.7, 4.45, 8.2]
    for x, (title, bullets) in zip(x_positions, columns):
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.45), Inches(3.35), Inches(4.7))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = SKY
        panel.line.width = Pt(1.2)
        add_text(slide, Inches(x + 0.2), Inches(1.7), Inches(2.9), Inches(0.3), title, size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_bullets(slide, Inches(x + 0.18), Inches(2.15), Inches(2.95), Inches(3.55), bullets, font_size=16)

    add_text(
        slide,
        Inches(0.9),
        Inches(6.5),
        Inches(11.5),
        Inches(0.3),
        "Overall direction: preserve the simulator’s accessibility while increasing its analytical depth and instructional flexibility.",
        size=13,
        color=MUTED,
    )
    add_footer(slide, "Future work balances usability, realism, and educational impact")


def main() -> None:
    data = run_analysis()
    prs = Presentation()
    set_slide_size(prs)

    add_title_slide(prs, data)
    add_problem_slide(prs)
    add_methodology_slide(prs)
    add_design_slide(prs)
    add_results_slide(prs, data)
    add_interpretation_slide(prs, data)
    add_future_work_slide(prs)

    prs.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
